# -*- coding: utf-8 -*-
"""
    This pipeline compiles the Non-SRT transcripts and PPT text into two
    seperate JSON files. The schema for them is as follows:

Transcript Schema:

        [
            {
            "module_number": 0,
            "module_name":"",
            "file_name": "",
            "transcript":[
                 {
                     "title": "",
                     “data": ["", ...],
                 }, ...
             ],
        ]


PPT Schema:

        [
            {
            "module_number": 0,
            "module_name":"",
            "file_name": "",
            "ppt":[
                 {
                     "slide_number": 0,
                     "slide_title" : "",
                     "slide_text" : []
                 }, ...
             ],
        ]

"""

import os
import json
from docx import Document # pip install python-docx
from pptx import Presentation # pip install python-pptx


# Raw data location
loc_transcript = "data/Non-SRT Transcript/"
loc_ppt = "data/PPT/"


# --- Transcript --------------------------------------------------------------


# Processed JSON dict
data_transcript = []


# Create module level JSON structure for transcripts
for f in os.listdir(loc_transcript):
    file = f.split(".")
    if file[1] == "docx":
        file = file[0].split()
        data_transcript.append({
            "module_number": int(file[1]),
            "module_name": ' '.join(file[2:]),
            "file_name": f,
            "transcript": [],
            })


# Sort JSON for transcripts
data_transcript.sort(key=lambda x: x['module_number'])


# Check if the para is a title or not
def if_title(doctext):
    words = doctext.split(" ")
    w = words[len(words)-1]
    try:
        float(w)
        result = True
    except:
        result = False
    if (result):
        return True
    return False


# Extract transcript text
for d in data_transcript:
    doc = Document(loc_transcript + d["file_name"])
    # print(d["file_name"])

    # Doc is a word document object. We can extract paragraphs from this obj.
    # Check out the "python-docx" library documentation
    para_count  = 0 # Count of para in the doc obj
    transcript_data = [] # Data for every title in topic:transcript
    for t in doc.paragraphs:
        para_count += 1

        # Ignore if the para line is the heading, say 'Week 1 Module 1 Fundamentals of System Hardware'
        if para_count == 1:
            continue

        resTitle = if_title(t.text)
        # print(resTitle)

        # Append transcript title and data
        if (resTitle) and (not transcript_data):
            title = t.text

        elif not(resTitle):
            transcript_data.append(t.text)

        else:
            d['transcript'].append({
                "title": title,
                "data": transcript_data
            })
            title = t.text
            transcript_data = []



# --- PPT ---------------------------------------------------------------------


# Processed JSON dict
data_ppt = []


# Create module level JSON structure for PPTs
for f in os.listdir(loc_ppt):
    file = f.split(".")
    if file[1] == "pptx":
        file = file[0].split(' - ')
        data_ppt.append({
            "module_number": int(file[0]),
            "module_name": ' '.join(file[1:]),
            "file_name": f,
            "ppt": [],
            })


# Sort JSON for PPTs
data_ppt.sort(key=lambda x: x['module_number'])


# Extract PPT text
for d in data_ppt:
    ppt = Presentation(loc_ppt + d["file_name"])
    # print(d)
    for idx1, slide in enumerate(ppt.slides):
        slide_number = idx1
        slide_text = []
        slide_title = ""
        slide_text_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if slide_text_count == 0:
                    slide_title = shape.text
                    slide_text.append(shape.text)
                else:
                    slide_text.append(shape.text)
                slide_text_count += 1
                # print(shape.text)
        d["ppt"].append({
            "slide_number" : slide_number,
            "slide_title" : slide_title,
            "slide_text" : slide_text
        })


# --- Post-Processing ---------------------------------------------------------


# Print data
# print(data_transcript)
# print(data_ppt)


# Save JSON
with open("data/preprocessed/transcripts.json", "w") as f:
    json.dump(data_transcript, f, indent = 4)

with open("data/preprocessed/ppts.json", "w") as f:
    json.dump(data_ppt, f, indent = 4)